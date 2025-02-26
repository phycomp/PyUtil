from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob
import os
import codecs
from PIL import Image
import pytesseract
pytesseract.pytesseract.tesseract_cmd = '/usr/local/Cellar/tesseract/4.1.1/bin/tesseract'
from pytesseract import image_to_string

n=0
def write_image(shape):
    global n
    image = shape.image
    # get image 
    image_bytes = image.blob
    # assinging file name, e.g. 'image.jpg'
    image_filename = fname[:-5]+'{:03d}.{}'.format(n, image.ext)
    n += 1
    print(image_filename)
    os.chdir("directory_path/readpptx/images")
    with open(image_filename, 'wb') as fout:
        fout.write(image_bytes)
    os.chdir("directory_path/readpptx")    

def visitor(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            write_image(shape)

def iter_picture_shapes(prs1):
    for slide in prs1.slides:
        for shape in slide.shapes:
                visitor(shape)

with open("directory_path/MyFile.txt","a+") as fout:
  for each_file in glob.glob("directory_path/*.pptx"):
    fname = os.path.basename(each_file)
    fout.write("-------------------"+fname+"----------------------\n")
    prs = Presentation(each_file)
    print("---------------"+fname+"-------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                rndrCode(shape.text)
                fout.write(shape.text+"\n")
    iter_picture_shapes(prs)
