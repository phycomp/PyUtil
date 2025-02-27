from pptx import Presentation

def mkPPT():
  prs = Presentation()

# Slide 1: Title Slide
  slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
  title_1 = slide_1.shapes.title
  title_1.text = "Franchise Business Model for Beauty and Massage Market"
  subtitle_1 = slide_1.placeholders[1]
  subtitle_1.text = "Expanding through Franchise Opportunities"

# Slide 2: Key Elements of a Franchise Model
  slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
  title_2 = slide_2.shapes.title
  title_2.text = "Key Elements of the Franchise Model"

  content_2 = slide_2.shapes.placeholders[1].text_frame
  content_2.text = "1. Brand Recognition"
  p = content_2.add_paragraph()
  p.text = "2. Initial Investment & Fees"
  p = content_2.add_paragraph()
  p.text = "3. Training & Support"
  p = content_2.add_paragraph()
  p.text = "4. Exclusive Product Lines"
  p = content_2.add_paragraph()
  p.text = "5. Business Operations and Systems"
  p = content_2.add_paragraph()
  p.text = "6. Target Audience"
  p = content_2.add_paragraph()
  p.text = "7. Marketing and Promotions"

# Slide 3: Initial Investment & Fees
  slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
  title_3 = slide_3.shapes.title
  title_3.text = "Initial Investment & Fees"

  content_3 = slide_3.shapes.placeholders[1].text_frame
  content_3.text = "• Franchise Fee: Typically includes rights to use the brand and system"
  p = content_3.add_paragraph()
  p.text = "• Ongoing Royalties: Paid as a percentage of revenue"
  p = content_3.add_paragraph()
  p.text = "• Setup Costs: Includes design, equipment, and initial inventory"

# Slide 4: Training & Support
  slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
  title_4 = slide_4.shapes.title
  title_4.text = "Training & Support"

  content_4 = slide_4.shapes.placeholders[1].text_frame
  content_4.text = "• Initial training for services and customer management"
  p = content_4.add_paragraph()
  p.text = "• Ongoing support in marketing, operations, and product development"
  p = content_4.add_paragraph()
  p.text = "• Assistance with location selection and setup"

# Slide 5: Brand and Marketing Strategy
  slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
  title_5 = slide_5.shapes.title
  title_5.text = "Brand and Marketing Strategy"

  content_5 = slide_5.shapes.placeholders[1].text_frame
  content_5.text = "• Leverage national brand recognition to attract customers"
  p = content_5.add_paragraph()
  p.text = "• Franchisees receive marketing material and support for local promotion"
  p = content_5.add_paragraph()
  p.text = "• Collaboration in regional and national advertising campaigns"

# Slide 6: Conclusion - Benefits of Franchising in the Beauty and Massage Market
  slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
  title_6 = slide_6.shapes.title
  title_6.text = "Benefits of Franchising in the Beauty and Massage Market"

  content_6 = slide_6.shapes.placeholders[1].text_frame
  content_6.text = "• Access to a proven business model with lower risks"
  p = content_6.add_paragraph()
  p.text = "• Ongoing support, ensuring consistency in service quality"
  p = content_6.add_paragraph()
  p.text = "• Opportunity to enter a growing wellness industry"

# Save the presentation
  ppt_path = "beautyMassageFranchise.pptx"
  prs.save(ppt_path)

  return ppt_path
