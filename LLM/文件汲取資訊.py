from streamlit import text_area, button, subheader, title, file_uploader as flUpldr
import pandas as pd
import pytesseract
import pdf2image
import pptx
from PIL import Image
import io
import textract
import spacy
from transformers import pipeline
from stUtil import rndrCode

class DocumentExtractor:
    def __init__(self):
      self.nlp = spacy.load("zh_core_web_sm") # 載入自然語言處理模型
      self.classifier = pipeline("zero-shot-classification")

    def extract_text(self, uploaded_file):
        """根據文件類型提取文字"""
        filename = uploaded_file.name.lower()
        if filename.endswith('.pdf'):
          return self._extract_pdf(uploaded_file)# PDF處理
        elif filename.endswith('.pptx'):
          return self._extract_pptx(uploaded_file)# PPT處理

        elif filename.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff')):
          return self._extract_image(uploaded_file)# 圖片處理

        else:
          return self._extract_text_file(uploaded_file) # 其他文字文件

    def _extract_pdf(self, file):
        """PDF文字提取"""
        images = pdf2image.convert_from_bytes(file.getvalue())
        texts = [pytesseract.image_to_string(img, lang='chi_tra') for img in images]
        return " ".join(texts)

    def _extract_pptx(self, file):
        """PPT文字提取"""
        prs = pptx.Presentation(file)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return " ".join(texts)

    def _extract_image(self, file):
        """圖片文字提取"""
        image = Image.open(file)
        return pytesseract.image_to_string(image, lang='chi_tra')

    def _extract_text_file(self, file):
        """通用文字檔提取"""
        return textract.process(file.name).decode('utf-8')

    def analyze_text(self, text): """文本分析"""
        doc = self.nlp(text)
        entities = [(ent.text, ent.label_) for ent in doc.ents] # 實體識別

        keywords = [token.text for token in doc if token.pos_ in ['NOUN', 'PROPN']] # 關鍵詞提取

        categories = ['產品介紹', '技術文檔', '行銷材料', '研究報告'] # AI分類
        classification = self.classifier(text, categories)

        return {
            'entities': entities,
            'keywords': keywords,
            'top_category': classification['labels'][0],
            'confidence': classification['scores'][0]
        }

def main():
    st.title("多格式文件智能提取系統")
    extractor = DocumentExtractor()
    uploaded_file = flUpldr("上傳文件", type=['pdf', 'pptx', 'png', 'jpg', 'txt'])

    if uploaded_file:
      extracted_text = extractor.extract_text(uploaded_file) # 文字提取
      text_area("提取文字", extracted_text, height=200)

      if button("智能分析"): # 文本分析
        analysis = extractor.analyze_text(extracted_text)
        subheader("分析結果")
        rndrCode("實體識別:", analysis['entities'])
        rndrCode("關鍵詞:", analysis['keywords'])
        rndrCode(f"文檔類型: {analysis['top_category']} (信心度: {analysis['confidence']:.2%})")

if __name__ == "__main__":
    main()

# 依賴安裝
"""
pip install streamlit pytesseract pdf2image python-pptx textract spacy transformers
python -m spacy download zh_core_web_sm
"""
