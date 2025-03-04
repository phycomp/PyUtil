from pptx import Presentation
from PIL import Image
from stUtil import rndrCode
import os

# 加載 PowerPoint 檔案
from sys import argv
ppt_file = argv[1]
prs = Presentation(ppt_file)

# 設置輸出資料夾
output_folder = 'output_images'
if not os.path.exists(output_folder):
    os.makedirs(output_folder)

# 逐頁保存為 PNG
for i, slide in enumerate(prs.slides):
    image_path = os.path.join(output_folder, f'slide_{i+1}.png')
    slide_img = slide.shapes  # 獲取幻燈片中的圖片
    # 這裡可選擇你如何獲取幻燈片的圖片並轉換成 PNG (如使用 PIL 進行操作)
    # 儲存 PNG (假設你有進行圖片轉換)
    slide_img.save(image_path)

rndrCode("轉換完成!")
